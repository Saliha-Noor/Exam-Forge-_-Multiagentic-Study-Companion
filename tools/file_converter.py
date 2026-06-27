import os


def convert_pptx(filepath):
    from pptx import Presentation
    print(f"  Converting PPTX: {os.path.basename(filepath)}")
    prs = Presentation(filepath)
    text = ""
    for i, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"
        if slide_text:
            text += f"\n--- Slide {i} ---\n{slide_text}"
    return text


def convert_docx(filepath):
    from docx import Document
    print(f"  Converting DOCX: {os.path.basename(filepath)}")
    doc = Document(filepath)
    text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text.strip() + "\n"
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                text += row_text + "\n"
    return text


def convert_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext in (".pptx", ".ppt"):
        return convert_pptx(filepath)
    elif ext == ".docx":
        return convert_docx(filepath)
    elif ext == ".txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        print(f"  Skipping: {os.path.basename(filepath)}")
        return None